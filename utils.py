from pptx import Presentation
import markdown
from pptx.util import Inches


def creer_slide(prs, text):
    slide_layout = prs.slide_layouts[5]  # Choisir le layout pour le titre et le contenu

    # Créer un nouveau slide
    slide = prs.slides.add_slide(slide_layout)

    # Ajouter le titre au slide
    title = slide.shapes.title
    title.text = "Points positifs et négatifs concernant l'hôtel"

    # Ajouter le contenu formaté au slide
    left = top = width = height = Inches(2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = text
